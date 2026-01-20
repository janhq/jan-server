#!/usr/bin/env python3
"""
render_deck.py
Render a DeckSpec-like JSON to PPTX using python-pptx.

Supported:
- theme.canvas.size: WIDE_16x9, STANDARD_4x3, or CUSTOM (inches)
- masters.globalElements applied to each slide via layout.masterId
- components referenced by slide.useComponents
- element types: text, shape, table, chart, image (file & URL), group (flatten)
- chart: bar|line|pie|area, dataset kind: series
- text: rich text runs, bullets, alignment
- shape: fill (solid/gradient), stroke, rotation, opacity
- speakerNotes

Usage:
  python render_deck.py input.json output.pptx
"""
# CRITICAL: Print immediately to confirm module loading starts
import sys
print(">>> MODULE LOAD: render_deck.py starting", file=sys.stderr, flush=True)

import json
import re
import logging
import tempfile
import io
from pathlib import Path
from urllib.parse import urlparse

print(">>> MODULE LOAD: stdlib imports OK", file=sys.stderr, flush=True)

print(">>> MODULE LOAD: stdlib imports OK", file=sys.stderr, flush=True)

try:
    import requests
    HAS_REQUESTS = True
    print(">>> MODULE LOAD: requests available", file=sys.stderr, flush=True)
except ImportError:
    HAS_REQUESTS = False
    print(">>> MODULE LOAD: requests NOT available", file=sys.stderr, flush=True)

print(">>> MODULE LOAD: importing pptx...", file=sys.stderr, flush=True)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
print(">>> MODULE LOAD: pptx imports OK", file=sys.stderr, flush=True)

print(">>> MODULE LOAD: pptx imports OK", file=sys.stderr, flush=True)

# ----------------------------
# Logging setup
# ----------------------------
print(">>> MODULE LOAD: configuring logging", file=sys.stderr, flush=True)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    stream=sys.stderr,  # Explicitly write to stderr for subprocess capture
    force=True  # Override any existing config
)
logger = logging.getLogger(__name__)
logger.info("Logger initialized")
print(">>> MODULE LOAD: logging configured", file=sys.stderr, flush=True)

# ----------------------------
# Constants
# ----------------------------
WIDE_16x9_DIMS = (13.333, 7.5)
STANDARD_4x3_DIMS = (10.0, 7.5)
EMU_PER_INCH = 914400  # PowerPoint internal units

# ----------------------------
# Helpers
# ----------------------------
def hex_to_rgb(hex_color: str):
    if not hex_color:
        return None
    h = hex_color.strip()
    if h.startswith("#"):
        h = h[1:]
    if len(h) == 8:  # ignore alpha
        h = h[:6]
    if len(h) != 6:
        return None
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    return RGBColor(r, g, b)

def emu_in(x_in: float):
    return Inches(float(x_in))

# Global unit setting - will be set by render()
_current_unit = "in"

def set_unit(unit: str):
    global _current_unit
    _current_unit = unit

def convert_to_inches(val: float) -> float:
    """Convert value from current unit to inches (no scaling)."""
    if _current_unit == "pt":
        return val / 72.0
    elif _current_unit == "cm":
        return val / 2.54
    return val  # already inches

def rect_in(rect: dict):
    """Convert rect coordinates to PowerPoint Inches based on current unit."""
    x = convert_to_inches(float(rect.get("x", 0)))
    y = convert_to_inches(float(rect.get("y", 0)))
    w = convert_to_inches(float(rect.get("w", 0)))
    h = convert_to_inches(float(rect.get("h", 0)))
    return (emu_in(x), emu_in(y), emu_in(w), emu_in(h))

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY
}
VALIGN_MAP = {
    "top": MSO_VERTICAL_ANCHOR.TOP,
    "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
    "bottom": MSO_VERTICAL_ANCHOR.BOTTOM
}



VAR_PATTERN = re.compile(r"\{\{\s*([a-zA-Z0-9_.-]+)\s*\}\}")

def substitute_vars(s: str, ctx: dict, page_num: int, total_pages: int) -> str:
    """Replace {page}, {total_pages}, and {{var}} placeholders using ctx."""
    if s is None:
        return ""
    out = str(s)
    out = out.replace("{page}", str(page_num))
    out = out.replace("{total_pages}", str(total_pages)).replace("{totalPages}", str(total_pages))

    def _lookup(key: str):
        if key in ("page", "pageNumber"):
            return str(page_num)
        if key in ("total_pages", "totalPages", "totalPagesCount"):
            return str(total_pages)
        if not ctx:
            return None
        if key in ctx and ctx[key] is not None:
            return ctx[key]
        if "." in key:
            cur = ctx
            for part in key.split('.'):
                if isinstance(cur, dict) and part in cur:
                    cur = cur[part]
                else:
                    return None
            return cur
        return None

    def repl(m):
        key = m.group(1)
        val = _lookup(key)
        return str(val) if val is not None else m.group(0)

    return VAR_PATTERN.sub(repl, out)
def safe_contains(rect, safe):
    """Check if rect is within safe margins (all values in inches after conversion)."""
    x = convert_to_inches(float(rect.get("x", 0)))
    y = convert_to_inches(float(rect.get("y", 0)))
    w = convert_to_inches(float(rect.get("w", 0)))
    h = convert_to_inches(float(rect.get("h", 0)))
    return (x >= safe["left"] and
            y >= safe["top"] and
            x + w <= (safe["slide_w"] - safe["right"]) and
            y + h <= (safe["slide_h"] - safe["bottom"]))

def set_fill(shape, fill_spec: dict):
    """Apply fill to a shape - supports solid and gradient fills."""
    if not fill_spec or fill_spec.get("type") in (None, "none"):
        return
    ftype = fill_spec.get("type")
    if ftype == "solid":
        color = fill_spec.get("color")
        rgb = hex_to_rgb(color)
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb
    elif ftype == "gradient":
        # python-pptx has limited gradient support, use first/last stops as approximation
        gradient = fill_spec.get("gradient", {})
        stops = gradient.get("stops", [])
        if stops:
            # Use first color as solid fill (gradient approximation)
            first_color = stops[0].get("color")
            rgb = hex_to_rgb(first_color)
            if rgb:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb
            logger.debug(f"Gradient fill approximated with solid color: {first_color}")

def set_line(shape, stroke_spec: dict):
    if not stroke_spec:
        return
    line = shape.line
    color = stroke_spec.get("color")
    width = stroke_spec.get("width")
    if color:
        rgb = hex_to_rgb(color)
        if rgb:
            line.color.rgb = rgb
    if width is not None:
        line.width = Pt(float(width))



def _remove_effects(shape):
    """Remove effect lists (shadow/glow) from a shape/picture/chart."""
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            return
        for tag in ('a:effectLst', 'a:effectDag'):
            eff = spPr.find(qn(tag))
            if eff is not None:
                spPr.remove(eff)
    except Exception:
        pass


def disable_shadow(shape):
    """Best-effort: disable any drop shadow (cleaner, flatter templates)."""
    try:
        sh = getattr(shape, 'shadow', None)
        if sh is not None:
            try:
                sh.inherit = False
            except Exception:
                pass
            try:
                sh.visible = False
            except Exception:
                pass
    except Exception:
        pass

    # Also remove effect lists in underlying XML (more reliable than the API)
    _remove_effects(shape)

def _enable_bullet(p, bullet_spec: dict):
    """Enable bullets on paragraph, with optional indent/hanging (in pt)."""
    try:
        from lxml import etree
        p.level = 0
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.find(qn('a:buChar'))
        if buChar is None:
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '•')

        # Indentation values in DrawingML are in EMU
        indent = bullet_spec.get('indent')
        hanging = bullet_spec.get('hanging')
        if indent is not None:
            marL = int((float(indent) / 72.0) * EMU_PER_INCH)
            pPr.set('marL', str(marL))
        if hanging is not None:
            ind = -int((float(hanging) / 72.0) * EMU_PER_INCH)
            pPr.set('indent', str(ind))
    except Exception:
        pass


def _extract_text_obj(el: dict) -> dict:
    text_obj = el.get("text")
    if isinstance(text_obj, dict):
        return text_obj

    content = el.get("content")
    if isinstance(content, dict):
        # Possible shapes:
        # {"text": {"content": "...", "style": {...}}}
        # {"text": "...", "style": {...}}
        # {"content": "...", "style": {...}}
        nested = content.get("text")
        if isinstance(nested, dict):
            return nested
        if isinstance(nested, str):
            return {"content": nested, "style": content.get("style", {})}
        raw = content.get("content")
        if isinstance(raw, str):
            return {"content": raw, "style": content.get("style", {})}
        if "style" in content:
            return {"content": "", "style": content.get("style", {})}

    if isinstance(content, str):
        return {"content": content}

    return {}

def add_text(slide, el, theme_defaults_text, slide_number: int, total_pages: int, text_ctx: dict):
    """Add text element with support for rich text runs, bullets, and auto-fit."""
    left, top, width, height = rect_in(el["rect"])
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True

    text_obj = _extract_text_obj(el)
    content = text_obj.get("content", "") or ""

    # If content uses "|" as newline delimiter, normalize it (common in prompts)
    if "|" in content and "\n" not in content:
        content = content.replace("|", "\n")

    # Substitute template variables for simple content (avoid breaking run offsets)
    runs = text_obj.get("runs", []) or []
    if not runs:
        content = substitute_vars(content, text_ctx, slide_number, total_pages)

    style = text_obj.get("style") or {}

    # Auto-fit controls (schema field text.autoFit)
    auto_fit = text_obj.get("autoFit") or "shrink"
    try:
        if auto_fit == "shrink":
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        elif auto_fit == "resizeBox":
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        else:
            tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass

    bullet_spec = style.get("bullet") or {}
    bullet_enabled = bool(bullet_spec.get("enabled", False))

    lines = content.split("\n") if content else [""]

    # Heuristic: if content already looks like a bulleted list, enable bullets
    if not bullet_enabled:
        markers = ("•", "-", "–", "—", "*")
        if any((ln.strip() and ln.lstrip().startswith(markers)) for ln in lines):
            bullet_enabled = True

    def _strip_bullet_marker(ln: str) -> str:
        t = ln.strip()
        for m in ("•", "-", "–", "—", "*"):
            if t.startswith(m):
                return t[len(m):].lstrip()
        return t

    char_offset = 0  # Track position in original content for runs

    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()

        line_start = char_offset
        line_end = char_offset + len(line)
        line_runs = [r for r in runs if r.get("start", 0) < line_end and r.get("end", 0) > line_start]

        if line_runs:
            # Rich text runs are applied verbatim; bullet-marker stripping is not safe here.
            _apply_rich_text_line(p, line, line_runs, line_start, style, theme_defaults_text)
        else:
            txt = _strip_bullet_marker(line) if bullet_enabled else line.strip()
            p.text = txt
            _apply_paragraph_style(p, style, theme_defaults_text)

        if bullet_enabled and (p.text or "").strip():
            _enable_bullet(p, bullet_spec)

        align = style.get("align")
        if align in ALIGN_MAP:
            p.alignment = ALIGN_MAP[align]

        char_offset = line_end + 1

    valign = style.get("valign")
    if valign in VALIGN_MAP:
        tf.vertical_anchor = VALIGN_MAP[valign]

    link = el.get("link", {}) or {}
    url = link.get("url")
    if url:
        p0 = tf.paragraphs[0]
        if p0.runs:
            p0.runs[0].hyperlink.address = url

    rotation = el.get("rotation")
    if rotation:
        tb.rotation = float(rotation)

    disable_shadow(tb)
    return tb





def _apply_paragraph_style(p, style: dict, theme_defaults_text: dict):
    """Apply style to all text in a paragraph."""
    font = p.font
    family = style.get("fontFamily") or theme_defaults_text.get("fontFamily")
    if family:
        font.name = family
    size = style.get("fontSize") or theme_defaults_text.get("fontSize")
    if size:
        font.size = Pt(float(size))
    if "bold" in style:
        font.bold = bool(style["bold"])
    if "italic" in style:
        font.italic = bool(style["italic"])
    if "underline" in style:
        font.underline = bool(style["underline"])
    color = style.get("color") or theme_defaults_text.get("color")
    rgb = hex_to_rgb(color) if color else None
    if rgb:
        font.color.rgb = rgb

    # Line height (multiple) if provided
    lh = style.get("lineHeight")
    if lh is None:
        lh = theme_defaults_text.get("lineHeight")
    if isinstance(lh, (int, float)):
        try:
            p.line_spacing = float(lh)
        except Exception:
            pass



def _apply_rich_text_line(p, line: str, runs: list, line_start: int, base_style: dict, theme_defaults_text: dict):
    """Apply rich text runs to a paragraph line."""
    # Clear existing text, we'll add runs manually
    p.clear()
    
    # Sort runs by start position
    runs = sorted(runs, key=lambda r: r.get("start", 0))
    
    current_pos = 0
    for run_spec in runs:
        run_start = max(0, run_spec.get("start", 0) - line_start)
        run_end = min(len(line), run_spec.get("end", 0) - line_start)
        
        if run_start > len(line) or run_end < 0:
            continue
            
        # Add text before this run with base style
        if run_start > current_pos:
            r = p.add_run()
            r.text = line[current_pos:run_start]
            _apply_run_style(r, base_style, theme_defaults_text)
        
        # Add the styled run
        if run_end > run_start:
            r = p.add_run()
            r.text = line[run_start:run_end]
            merged_style = {**base_style, **(run_spec.get("style", {}))}
            _apply_run_style(r, merged_style, theme_defaults_text)
            current_pos = run_end
    
    # Add remaining text after last run
    if current_pos < len(line):
        r = p.add_run()
        r.text = line[current_pos:]
        _apply_run_style(r, base_style, theme_defaults_text)


def _apply_run_style(run, style: dict, theme_defaults_text: dict):
    """Apply style to a text run."""
    font = run.font
    family = style.get("fontFamily") or theme_defaults_text.get("fontFamily")
    if family:
        font.name = family
    size = style.get("fontSize") or theme_defaults_text.get("fontSize")
    if size:
        font.size = Pt(float(size))
    if "bold" in style:
        font.bold = bool(style["bold"])
    if "italic" in style:
        font.italic = bool(style["italic"])
    if "underline" in style:
        font.underline = bool(style["underline"])
    color = style.get("color") or theme_defaults_text.get("color")
    rgb = hex_to_rgb(color) if color else None
    if rgb:
        font.color.rgb = rgb


def add_shape(slide, el):
    """Add shape element with rotation and opacity support.

    Guardrails:
    - No rounded rectangles or decorative circles: anything not in the allowed set is rendered as a plain rectangle.
    - Shadows are removed.
    """
    left, top, width, height = rect_in(el["rect"])
    s = el.get("shape", {})
    kind = (s.get("kind") or "rect")
    kind = str(kind)
    style = s.get("style", {}) or {}

    if kind in ("line", "arrow"):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top + height)
        set_line(conn, (style.get("stroke") or {}))
        if kind == "arrow":
            conn.line.end_arrowhead = True
        disable_shadow(conn)
        return conn

    allowed = {"rect", "triangle", "diamond"}
    if kind not in allowed:
        kind = "rect"

    SHAPE_MAP = {
        "rect": MSO_SHAPE.RECTANGLE,
        # Any rounded/oval requests are forced to RECTANGLE by design
        "roundRect": MSO_SHAPE.RECTANGLE,
        "roundedRect": MSO_SHAPE.RECTANGLE,
        "pill": MSO_SHAPE.RECTANGLE,
        "ellipse": MSO_SHAPE.RECTANGLE,
        "circle": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.RECTANGLE,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
    }
    shape_type = SHAPE_MAP.get(kind, MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_fill(shp, style.get("fill"))
    set_line(shp, (style.get("stroke") or {}))
    disable_shadow(shp)

    rotation = el.get("rotation")
    if rotation:
        shp.rotation = float(rotation)

    return shp


def add_table(slide, el, theme_defaults_text, slide_number: int):
    left, top, width, height = rect_in(el["rect"])
    t = el.get("table", {})
    cols = t.get("columns", [])
    rows = t.get("rows", [])
    nrows = 1 + len(rows)
    ncols = max(1, len(cols))
    table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = table_shape.table

    for c in range(ncols):
        hdr = cols[c] if c < len(cols) else ""
        table.cell(0, c).text = hdr

    for r_i, row in enumerate(rows, start=1):
        for c in range(ncols):
            val = row[c] if c < len(row) else ""
            table.cell(r_i, c).text = str(val)

    style = t.get("style", {})
    hdr_style = style.get("headerTextStyle", {})
    cell_style = style.get("cellTextStyle", {})

    def style_cell(cell, st):
        tf = cell.text_frame
        p = tf.paragraphs[0]
        font = p.font
        family = st.get("fontFamily") or theme_defaults_text.get("fontFamily")
        if family:
            font.name = family
        size = st.get("fontSize") or theme_defaults_text.get("fontSize")
        if size:
            font.size = Pt(float(size))
        color = st.get("color") or theme_defaults_text.get("color")
        rgb = hex_to_rgb(color) if color else None
        if rgb:
            font.color.rgb = rgb
        if "bold" in st:
            font.bold = bool(st["bold"])

    for c in range(ncols):
        style_cell(table.cell(0, c), hdr_style)
    for r in range(1, nrows):
        for c in range(ncols):
            style_cell(table.cell(r, c), cell_style)

    return table_shape

def _extract_chart_obj(el: dict) -> dict:
    chart_obj = el.get("chart")
    if isinstance(chart_obj, dict):
        return chart_obj

    content = el.get("content")
    if isinstance(content, dict):
        nested = content.get("chart")
        if isinstance(nested, dict):
            return nested
        if "chartType" in content or "datasetRef" in content:
            return content

    return {}


def add_chart(slide, el, datasets_by_id):
    """Add chart element with support for bar, line, pie, and area charts."""
    left, top, width, height = rect_in(el["rect"])
    c = _extract_chart_obj(el)
    chart_type = c.get("chartType", "bar")
    ds_id = c.get("datasetRef")
    ds = datasets_by_id.get(ds_id) if ds_id else None
    
    # If no dataset found, create placeholder data
    if not ds:
        logger.warning(f"Chart datasetRef '{ds_id}' not found, using placeholder data")
        categories = ["Category 1", "Category 2", "Category 3"]
        series = [{"name": "Placeholder", "values": [10, 20, 30]}]
    else:
        data = ds.get("data", {})
        categories = data.get("labels", ["A", "B", "C"])
        series = data.get("series", [{"name": "Data", "values": [1, 2, 3]}])
        # Ensure we have valid data
        if not categories:
            categories = ["A", "B", "C"]
        if not series:
            series = [{"name": "Data", "values": [1, 2, 3]}]

    chart_data = ChartData()
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s.get("name", "Series"), s.get("values", []))

    # Map chart types including area
    CHART_TYPE_MAP = {
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "area": XL_CHART_TYPE.AREA,
    }
    xl_type = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)

    chart_shape = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = chart_shape.chart

    title = c.get("title")
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title

    style = c.get("style", {})
    chart.has_legend = bool(style.get("showLegend", True))

    return chart_shape

def _extract_image_obj(el: dict) -> dict:
    image_obj = el.get("image")
    if isinstance(image_obj, dict):
        return image_obj

    content = el.get("content")
    if isinstance(content, dict):
        nested = content.get("image")
        if isinstance(nested, dict):
            return nested
        return content

    return {}


def add_image(slide, el, assets_by_id):
    """Add image element with support for file and URL sources."""
    left, top, width, height = rect_in(el["rect"])
    img = _extract_image_obj(el)
    ref = img.get("ref")
    asset = assets_by_id.get(ref)
    if not asset:
        fallback = None
        if ref:
            for key, value in assets_by_id.items():
                if isinstance(key, str) and (key.startswith(ref) or ref in key):
                    fallback = value
                    logger.warning(f"Image ref '{ref}' matched asset id '{key}'.")
                    break
        if not fallback:
            logger.warning(f"Image ref not found in assets: {ref}. Skipping image.")
            return None
        asset = fallback
    
    src = asset.get("source", {})
    source_type = src.get("type")
    
    image_stream = None
    
    if source_type == "url":
        # Fetch image from URL
        url = src.get("url")
        if not url:
            raise ValueError(f"Image source type is 'url' but no URL provided for ref: {ref}")
        
        if not HAS_REQUESTS:
            raise ImportError("'requests' library is required for URL images. Install with: pip install requests")
        
        logger.info(f"Fetching image from URL: {url}")
        try:
            response = requests.get(url, timeout=30, headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            })
            response.raise_for_status()
            image_stream = io.BytesIO(response.content)
            logger.info(f"Successfully fetched image: {len(response.content)} bytes")
        except requests.RequestException as e:
            raise ValueError(f"Failed to fetch image from URL {url}: {e}")
    
    elif source_type == "file":
        raw_path_str = src.get("filePath")
        if not raw_path_str:
            raise ValueError(f"Image source type is 'file' but no filePath provided for ref: {ref}")
        
        path = Path(raw_path_str)
        
        # Fallback to local file if absolute path not found
        if not path.exists():
            local_path = Path.cwd() / path.name
            if local_path.exists():
                path = local_path
            else:
                # Also try relative to script if CWD is different
                script_dir_path = Path(__file__).parent / path.name
                if script_dir_path.exists():
                    path = script_dir_path
                else:
                    raise FileNotFoundError(f"Image file not found: {path} (also checked {local_path})")
        
        image_stream = str(path)
    
    elif source_type == "generated":
        # Placeholder for AI-generated images - log warning and skip
        prompt = src.get("prompt", "")
        logger.warning(f"Generated image type not yet supported. Prompt: {prompt[:50]}...")
        return None
    
    else:
        raise ValueError(f"Unsupported image source type: {source_type}. Use 'file' or 'url'.")
    
    if image_stream is None:
        return None
    
    pic = slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
    
    # Apply rotation if specified
    rotation = el.get("rotation")
    if rotation:
        pic.rotation = float(rotation)
    
    return pic

def flatten_elements(elements, dx: float = 0.0, dy: float = 0.0, dz: int = 0):
    """Flatten nested group structures, applying group rect offsets to children.

    This fixes common overlap bugs where children coordinates are relative to the group.
    """
    flat = []
    for el in elements or []:
        if not isinstance(el, dict):
            continue

        if el.get("type") == "group":
            g = el.get("group") or {}
            children = g.get("children") or []
            g_rect = el.get("rect") or g.get("rect") or {}
            try:
                gx = float(g_rect.get("x", 0))
                gy = float(g_rect.get("y", 0))
            except (TypeError, ValueError):
                gx, gy = 0.0, 0.0

            try:
                gz = int(el.get("zIndex", 0))
            except (TypeError, ValueError):
                gz = 0

            flat.extend(flatten_elements(children, dx + gx, dy + gy, dz + gz))
            continue

        # Non-group: copy and offset rect
        out = dict(el)
        r = dict(out.get("rect") or {})
        if r:
            try:
                r["x"] = float(r.get("x", 0)) + dx
                r["y"] = float(r.get("y", 0)) + dy
            except (TypeError, ValueError):
                pass
            out["rect"] = r

        try:
            out["zIndex"] = int(out.get("zIndex", 0)) + dz
        except (TypeError, ValueError):
            out["zIndex"] = dz

        flat.append(out)

    return flat

def resolve_top_text_overlaps(slide, region_in: float = 1.75, pad_pt: float = 10.0):
    """Very small, targeted overlap fix:

    If two text boxes in the top region overlap significantly in X, push the lower one down.
    This addresses the common title/subtitle overlap when titles wrap unexpectedly.
    """
    try:
        region = int(region_in * EMU_PER_INCH)
        pad = int((pad_pt / 72.0) * EMU_PER_INCH)

        texts = [sh for sh in slide.shapes if getattr(sh, 'has_text_frame', False)]
        texts.sort(key=lambda sh: int(sh.top))

        # Two passes to reduce cascading overlaps
        for _ in range(2):
            for i in range(len(texts) - 1):
                upper = texts[i]
                lower = texts[i + 1]
                if int(upper.top) > region or int(lower.top) > region:
                    continue

                ux1, uy1 = int(upper.left), int(upper.top)
                ux2, uy2 = ux1 + int(upper.width), uy1 + int(upper.height)
                lx1, ly1 = int(lower.left), int(lower.top)
                lx2, ly2 = lx1 + int(lower.width), ly1 + int(lower.height)

                x_overlap = min(ux2, lx2) - max(ux1, lx1)
                if x_overlap <= 0:
                    continue
                if x_overlap < 0.4 * min(int(upper.width), int(lower.width)):
                    continue

                # If Y-overlap, shift lower down
                if ly1 < uy2 and uy1 < ly2:
                    new_top = uy2 + pad
                    if new_top > ly1:
                        lower.top = new_top
    except Exception:
        pass


def render(deck: dict, out_path: Path):
    prs = Presentation()

    theme = deck.get("theme", {})
    canvas = theme.get("canvas", {})
    size = canvas.get("size", "WIDE_16x9")
    unit = canvas.get("unit", "in")
    
    # Map common canvas size strings to dimensions (in inches)
    SIZE_MAP = {
        "WIDE_16x9": WIDE_16x9_DIMS,
        "1920x1080": WIDE_16x9_DIMS,
        "1280x720": WIDE_16x9_DIMS,
        "STANDARD_4x3": STANDARD_4x3_DIMS,
        "1024x768": STANDARD_4x3_DIMS,
    }
    
    if size in SIZE_MAP:
        slide_w, slide_h = SIZE_MAP[size]
    elif size == "CUSTOM":
        cs = canvas.get("customSize", {}) or {}
        slide_w, slide_h = float(cs.get("width", 13.333)), float(cs.get("height", 7.5))
    else:
        logger.warning(f"Unknown canvas size '{size}', defaulting to WIDE_16x9")
        slide_w, slide_h = WIDE_16x9_DIMS

    # Set global unit for rect_in conversions (no scaling - use coordinates as-is)
    set_unit(unit)
    logger.info(f"Canvas: {size}, unit: {unit}, slide size: {slide_w}x{slide_h} inches")
    
    # Convert safe margins to inches based on unit
    def to_inches(val, src_unit):
        if src_unit == "pt":
            return val / 72.0
        elif src_unit == "cm":
            return val / 2.54
        return val  # already inches

    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)

    safe = canvas.get("safeMargins", {"top": 0.5, "right": 0.5, "bottom": 0.5, "left": 0.5})
    safe_ctx = {
        "top": to_inches(safe.get("top", 0.5), unit),
        "right": to_inches(safe.get("right", 0.5), unit),
        "bottom": to_inches(safe.get("bottom", 0.5), unit),
        "left": to_inches(safe.get("left", 0.5), unit),
    }
    safe_ctx["slide_w"] = slide_w
    safe_ctx["slide_h"] = slide_h
    safe_ctx["unit"] = unit

    defaults = (theme.get("defaults", {}) or {})
    theme_text_default = defaults.get("textStyle", {}) or {}

    masters = {m["id"]: m for m in deck.get("masters", []) if isinstance(m, dict) and "id" in m}
    layouts = {l["id"]: l for l in deck.get("layouts", []) if isinstance(l, dict) and "id" in l}
    components = {c["id"]: c for c in deck.get("components", []) if isinstance(c, dict) and "id" in c}

    assets = deck.get("assets", {}) or {}
    images = assets.get("images", []) or []
    assets_by_id = {img["id"]: img for img in images if isinstance(img, dict) and "id" in img}

    data = deck.get("data", {}) or {}
    datasets = data.get("datasets", []) or []
    datasets_by_id = {ds["id"]: ds for ds in datasets if isinstance(ds, dict) and "id" in ds}

    slides = sorted(deck.get("slides", []), key=lambda s: s.get("order", 0))

    # Context for template variable substitution in text elements
    text_ctx = {}
    meta = deck.get("metadata") or {}
    if isinstance(meta, dict):
        for k, v in meta.items():
            if v is not None:
                text_ctx[k] = v
    constraints_text = None
    if isinstance(meta, dict):
        constraints_text = meta.get("constraints")
    if isinstance(constraints_text, str):
        for line in constraints_text.splitlines():
            if ":" in line:
                k, v = line.split(":", 1)
                k = k.strip()
                v = v.strip()
                if k and v and k not in text_ctx:
                    text_ctx[k] = v
    
    # Get blank layout safely
    try:
        blank_layout = prs.slide_layouts[6]  # blank
    except IndexError:
        logger.warning("Slide layout index 6 not found, using last available layout")
        blank_layout = prs.slide_layouts[-1]

    logger.info(f"Rendering {len(slides)} slides at {slide_w}x{slide_h} inches")

    for idx, s in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)

        layout_id = s.get("layoutId")
        layout = layouts.get(layout_id)
        if not layout:
            logger.warning(f"Slide {s.get('id')} references unknown layoutId: {layout_id}. Rendering without master.")
        master = masters.get(layout.get("masterId")) if layout else None

        # Background
        bg = s.get("backgroundOverride") or (master.get("background") if master else None) or defaults.get("background")
        if bg and bg.get("type") == "solid" and bg.get("color"):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(bg["color"])

        all_elements = []
        if master:
            all_elements.extend(master.get("globalElements", []))
        for comp_id in s.get("useComponents", []) or []:
            comp = components.get(comp_id)
            if comp:
                all_elements.extend(comp.get("elements", []))
        all_elements.extend(s.get("elements", []))

        all_elements = flatten_elements(all_elements)
        all_elements.sort(key=lambda e: e.get("zIndex", 0))

        for el in all_elements:
            rect = el.get("rect", {})
            if rect and not safe_contains(rect, safe_ctx):
                logger.warning(f"Element {el.get('id')} on slide {idx} is outside safe margins: {rect}")
                # Continue anyway instead of raising - let user see the output

            etype = el.get("type")
            try:
                if etype == "text":
                    add_text(slide, el, theme_text_default, idx, len(slides), text_ctx)
                elif etype == "shape":
                    add_shape(slide, el)
                elif etype == "table":
                    add_table(slide, el, theme_text_default, idx)
                elif etype == "chart":
                    add_chart(slide, el, datasets_by_id)
                elif etype == "image":
                    add_image(slide, el, assets_by_id)
                else:
                    logger.warning(f"Unknown element type '{etype}' on slide {idx}, skipping")
            except Exception as e:
                logger.error(f"Error rendering element {el.get('id')} on slide {idx}: {e}")
                raise

        # Small overlap fix for top-of-slide text boxes
        resolve_top_text_overlaps(slide)

        notes = s.get("speakerNotes")
        if notes:
            ns = slide.notes_slide
            ns.notes_text_frame.text = notes

    # Validate that slides were actually added
    if len(prs.slides) == 0:
        raise ValueError("Cannot save presentation: no slides were added to the presentation")
    
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    logger.info(f"Successfully saved presentation to {out_path}")
    
    # Verify file was actually created
    if not out_path.exists():
        raise IOError(f"Failed to create output file: {out_path}")
    
    file_size = out_path.stat().st_size
    if file_size < 1000:
        raise ValueError(f"Generated PPTX is suspiciously small ({file_size} bytes) - likely corrupt")
    
    logger.info(f"Verified output file: {file_size} bytes")


def main():
    # Immediate output to confirm script execution
    print("=== render_deck.py starting ===")
    sys.stdout.flush()
    
    if len(sys.argv) < 3:
        print("Usage: python render_deck.py input.json output.pptx", file=sys.stderr)
        sys.exit(2)

    in_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2])
    
    print(f"Input: {in_path}")
    print(f"Output: {out_path}")
    sys.stdout.flush()

    if not in_path.exists():
        logger.error(f"Input file not found: {in_path}")
        sys.exit(1)

    try:
        deck = json.loads(in_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        error_msg = f"Invalid JSON in {in_path}: {e}"
        logger.error(error_msg)
        print(f"ERROR: {error_msg}", file=sys.stderr)
        sys.stderr.flush()
        sys.exit(1)
    
    # Validate deck structure
    if not isinstance(deck, dict):
        error_msg = f"Deck JSON is not an object: {type(deck)}"
        print(f"ERROR: {error_msg}", file=sys.stderr)
        sys.stderr.flush()
        sys.exit(1)
    
    slides = deck.get("slides", [])
    print(f"Deck loaded: {len(slides)} slides, version={deck.get('version', 'unknown')}")
    sys.stdout.flush()
    
    if not slides:
        error_msg = "Deck has no slides array or slides array is empty"
        print(f"ERROR: {error_msg}", file=sys.stderr)
        sys.stderr.flush()
        sys.exit(1)

    try:
        render(deck, out_path)
        success_msg = f"SUCCESS: Wrote {out_path} ({out_path.stat().st_size} bytes)"
        print(success_msg)
        sys.stdout.flush()  # Ensure output is captured by subprocess
        logger.info(success_msg)
    except Exception as e:
        logger.error(f"Rendering failed: {e}")
        # Print to stderr for visibility in subprocess output
        error_msg = f"ERROR: Rendering failed: {e}"
        print(error_msg, file=sys.stderr)
        sys.stderr.flush()  # Ensure error is captured
        import traceback
        traceback.print_exc(file=sys.stderr)
        sys.stderr.flush()
        sys.exit(1)


if __name__ == "__main__":
    print(">>> SCRIPT ENTRY: __main__ block starting", file=sys.stderr, flush=True)
    try:
        print(">>> SCRIPT ENTRY: calling main()", file=sys.stderr, flush=True)
        main()
        print(">>> SCRIPT ENTRY: main() completed successfully", file=sys.stderr, flush=True)
    except Exception as e:
        # Catch any uncaught exceptions and print them
        print(f">>> SCRIPT ENTRY: FATAL ERROR in main(): {e}", file=sys.stderr, flush=True)
        import traceback
        traceback.print_exc(file=sys.stderr)
        sys.stderr.flush()
        sys.exit(1)
